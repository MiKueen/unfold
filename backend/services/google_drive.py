import io
import re
import zipfile
import logging
import xml.etree.ElementTree as ET
from typing import Optional
import requests
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from google.oauth2.credentials import Credentials

logger = logging.getLogger(__name__)

GOOGLE_EXPORT_MAP = {
    "application/vnd.google-apps.document": "text/plain",
    "application/vnd.google-apps.spreadsheet": "text/csv",
    "application/vnd.google-apps.presentation": "text/plain",
}

SUPPORTED_MIME_PREFIXES = [
    "text/",
    "application/json",
    "application/xml",
    "application/javascript",
    "application/pdf",
    "application/vnd.google-apps.document",
    "application/vnd.google-apps.spreadsheet",
    "application/vnd.google-apps.presentation",
    # Uploaded Office files
    "application/vnd.openxmlformats-officedocument.spreadsheetml",        # .xlsx
    "application/vnd.openxmlformats-officedocument.wordprocessingml",     # .docx
    "application/vnd.openxmlformats-officedocument.presentationml",       # .pptx
    "application/vnd.ms-excel",                                           # .xls
    "application/msword",                                                 # .doc
    "application/vnd.ms-powerpoint",                                      # .ppt
]

MAX_CHARS = 100_000 # Increased for production rigor
MAX_FILES = 200


def is_supported(mime_type: str) -> bool:
    return any(mime_type.startswith(p) for p in SUPPORTED_MIME_PREFIXES)


def validate_access_token(access_token: str) -> bool:
    """Production Rigor: Validate the token against Google's API."""
    try:
        # Call tokeninfo to verify the token is valid and not expired
        resp = requests.get(f"https://www.googleapis.com/oauth2/v3/tokeninfo?access_token={access_token}", timeout=5)
        return resp.status_code == 200
    except Exception:
        return False


def create_drive_service(access_token: str):
    creds = Credentials(token=access_token)
    return build("drive", "v3", credentials=creds)


def extract_drive_id(url: str) -> tuple[Optional[str], str]:
    """
    Returns (id, type) where type is 'folder', 'file', or 'unknown'.
    'unknown' means the type needs to be resolved via the Drive API.
    """
    url = url.strip()

    # Raw ID (no URL scheme)
    if re.match(r"^[a-zA-Z0-9_-]{25,}$", url):
        return url, "unknown"

    # Folder URLs
    folder_match = re.search(r"/folders/([a-zA-Z0-9_-]+)", url)
    if folder_match:
        return folder_match.group(1), "folder"

    # Google Workspace and Drive file URLs
    file_patterns = [
        r"/document/d/([a-zA-Z0-9_-]+)",
        r"/spreadsheets/d/([a-zA-Z0-9_-]+)",
        r"/presentation/d/([a-zA-Z0-9_-]+)",
        r"/forms/d/([a-zA-Z0-9_-]+)",
        r"/file/d/([a-zA-Z0-9_-]+)",
        r"/drawings/d/([a-zA-Z0-9_-]+)",
    ]
    for pattern in file_patterns:
        m = re.search(pattern, url)
        if m:
            return m.group(1), "file"

    # ?id= or &id= (Drive open/uc links — type ambiguous)
    id_match = re.search(r"[?&]id=([a-zA-Z0-9_-]+)", url)
    if id_match:
        return id_match.group(1), "unknown"

    return None, "unknown"


# Keep old name as alias for any remaining internal callers
def extract_folder_id(url: str) -> Optional[str]:
    file_id, _ = extract_drive_id(url)
    return file_id


def get_single_file_as_folder(service, file_id: str) -> dict:
    """
    Fetches metadata for a single file and returns it wrapped in the same
    folder-list response structure that list_folder_recursive returns.
    """
    meta = service.files().get(
        fileId=file_id,
        fields="id,name,mimeType,size,modifiedTime,webViewLink",
        supportsAllDrives=True,
    ).execute()

    mime_type = meta.get("mimeType", "")
    if mime_type == "application/vnd.google-apps.folder":
        # Caller passed a folder URL via the ?id= path — delegate
        return list_folder_recursive(service, file_id)

    if not is_supported(mime_type):
        raise ValueError(
            f"File type '{mime_type}' is not supported. "
            "Supported types: Google Docs, Sheets, Slides, PDFs, text, and code files."
        )

    entry = {
        "id": meta["id"],
        "name": meta["name"],
        "mimeType": mime_type,
        "size": meta.get("size"),
        "modifiedTime": meta.get("modifiedTime"),
        "webViewLink": meta.get("webViewLink"),
        "path": meta["name"],
    }

    folder = {
        "id": file_id,
        "name": meta["name"],
        "files": [entry],
        "subfolders": [],
    }

    return {
        "folder": folder,
        "totalFiles": 1,
        "supportedFiles": 1,
        "flatFiles": [entry],
        "unsupportedCount": 0,
    }


def list_folder_recursive(service, folder_id: str, max_files: int = MAX_FILES) -> dict:
    flat_files = []
    unsupported = []

    def _list(parent_id: str, path: str) -> dict:
        folder = {"id": parent_id, "name": path.split("/")[-1] or "Root", "files": [], "subfolders": []}
        page_token = None

        while True:
            if len(flat_files) >= max_files:
                break

            params = {
                "q": f"'{parent_id}' in parents and trashed = false",
                "fields": "nextPageToken, files(id, name, mimeType, size, modifiedTime, webViewLink)",
                "pageSize": 100,
                "supportsAllDrives": True,
                "includeItemsFromAllDrives": True,
            }
            if page_token:
                params["pageToken"] = page_token

            res = service.files().list(**params).execute()
            files = res.get("files", [])
            page_token = res.get("nextPageToken")

            for f in files:
                if len(flat_files) >= max_files:
                    break
                entry = {
                    "id": f["id"],
                    "name": f["name"],
                    "mimeType": f["mimeType"],
                    "size": f.get("size"),
                    "modifiedTime": f.get("modifiedTime"),
                    "webViewLink": f.get("webViewLink"),
                    "path": path,
                }
                if f["mimeType"] == "application/vnd.google-apps.folder":
                    sub = _list(f["id"], f"{path}/{f['name']}")
                    sub["name"] = f["name"]
                    folder["subfolders"].append(sub)
                elif is_supported(f["mimeType"]):
                    folder["files"].append(entry)
                    flat_files.append(entry)
                else:
                    unsupported.append(entry)

            if not page_token:
                break

        return folder

    folder_name = "Folder"
    try:
        meta = service.files().get(fileId=folder_id, fields="name", supportsAllDrives=True).execute()
        folder_name = meta.get("name", "Folder")
    except Exception:
        pass

    tree = _list(folder_id, folder_name)
    tree["name"] = folder_name

    return {
        "folder": tree,
        "totalFiles": len(flat_files) + len(unsupported),
        "supportedFiles": len(flat_files),
        "flatFiles": flat_files,
        "unsupportedCount": len(unsupported),
        "is_limited": len(flat_files) >= max_files
    }


_XLSX_NS = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"
_DOCX_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"


def _parse_xlsx(buf: io.BytesIO) -> Optional[str]:
    """Parse xlsx using stdlib zipfile + xml. No external dependencies."""
    with zipfile.ZipFile(buf) as zf:
        names = zf.namelist()

        # Load shared string table
        shared = []
        if "xl/sharedStrings.xml" in names:
            with zf.open("xl/sharedStrings.xml") as f:
                tree = ET.parse(f)
            for si in tree.findall(f".//{{{_XLSX_NS}}}si"):
                texts = [t.text or "" for t in si.findall(f".//{{{_XLSX_NS}}}t")]
                shared.append("".join(texts))

        # Load workbook to get sheet names
        sheet_names: dict[str, str] = {}
        if "xl/workbook.xml" in names:
            with zf.open("xl/workbook.xml") as f:
                wb_tree = ET.parse(f)
            wb_ns = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"
            for sh in wb_tree.findall(f".//{{{wb_ns}}}sheet"):
                rid = sh.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id", "")
                sheet_names[rid] = sh.get("name", "Sheet")

        # Load relationships to map rId → sheet file
        rid_to_file: dict[str, str] = {}
        rel_path = "xl/_rels/workbook.xml.rels"
        if rel_path in names:
            with zf.open(rel_path) as f:
                rel_tree = ET.parse(f)
            for rel in rel_tree.findall(".//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship"):
                rid_to_file[rel.get("Id", "")] = "xl/" + rel.get("Target", "").lstrip("/")

        # Parse sheets in order
        sheet_files = sorted(n for n in names if n.startswith("xl/worksheets/sheet") and n.endswith(".xml"))
        rows_all: list[str] = []

        for sheet_file in sheet_files:
            # Find sheet name
            sheet_label = sheet_file
            for rid, path in rid_to_file.items():
                if path == sheet_file:
                    sheet_label = sheet_names.get(rid, sheet_file)
                    break
            rows_all.append(f"=== Sheet: {sheet_label} ===")

            with zf.open(sheet_file) as f:
                sh_tree = ET.parse(f)

            for row_el in sh_tree.findall(f".//{{{_XLSX_NS}}}row"):
                cells: list[str] = []
                for c in row_el.findall(f"{{{_XLSX_NS}}}c"):
                    t = c.get("t", "")
                    v_el = c.find(f"{{{_XLSX_NS}}}v")
                    is_el = c.find(f".//{{{_XLSX_NS}}}t")  # inlineStr
                    if t == "s" and v_el is not None:
                        idx = int(v_el.text)
                        cells.append(shared[idx] if idx < len(shared) else "")
                    elif t == "inlineStr" and is_el is not None:
                        cells.append(is_el.text or "")
                    elif v_el is not None:
                        cells.append(v_el.text or "")
                    else:
                        cells.append("")
                if any(c.strip() for c in cells):
                    rows_all.append("\t".join(cells))

        return "\n".join(rows_all) or None


def _parse_docx(buf: io.BytesIO) -> Optional[str]:
    """Parse docx using stdlib zipfile + xml."""
    with zipfile.ZipFile(buf) as zf:
        if "word/document.xml" not in zf.namelist():
            return None
        with zf.open("word/document.xml") as f:
            tree = ET.parse(f)
    lines = []
    for para in tree.findall(f".//{{{_DOCX_NS}}}p"):
        parts = [t.text or "" for t in para.findall(f".//{{{_DOCX_NS}}}t")]
        text = "".join(parts).strip()
        if text:
            lines.append(text)
    return "\n".join(lines) or None


def _parse_pptx(buf: io.BytesIO) -> Optional[str]:
    """Parse pptx using stdlib zipfile + xml."""
    PPTX_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
    with zipfile.ZipFile(buf) as zf:
        slide_files = sorted(n for n in zf.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml"))
        parts = []
        for i, slide_file in enumerate(slide_files, 1):
            with zf.open(slide_file) as f:
                tree = ET.parse(f)
            texts = [t.text.strip() for t in tree.findall(f".//{{{PPTX_NS}}}t") if t.text and t.text.strip()]
            if texts:
                parts.append(f"--- Slide {i} ---\n" + "\n".join(texts))
    return "\n\n".join(parts) or None


def _parse_office_file(buf: io.BytesIO, mime_type: str) -> Optional[str]:
    """Extract text from uploaded Office files using stdlib parsers (no external deps required)."""
    try:
        if "spreadsheetml" in mime_type or mime_type == "application/vnd.ms-excel":
            text = _parse_xlsx(buf)
            if text:
                logger.info(f"Parsed xlsx via stdlib: {len(text)} chars")
                return text[:MAX_CHARS]
            # Try openpyxl as fallback
            try:
                buf.seek(0)
                import openpyxl
                wb = openpyxl.load_workbook(buf, read_only=True, data_only=True)
                rows = []
                for sheet in wb.worksheets:
                    rows.append(f"=== Sheet: {sheet.title} ===")
                    for row in sheet.iter_rows(values_only=True):
                        cells = [str(c) if c is not None else "" for c in row]
                        if any(c.strip() for c in cells):
                            rows.append("\t".join(cells))
                wb.close()
                text = "\n".join(rows)
                if text.strip():
                    logger.info(f"Parsed xlsx via openpyxl: {len(text)} chars")
                    return text[:MAX_CHARS]
            except Exception as e:
                logger.warning(f"openpyxl also failed: {e}")
            return None

        if "wordprocessingml" in mime_type or mime_type == "application/msword":
            text = _parse_docx(buf)
            if text:
                logger.info(f"Parsed docx via stdlib: {len(text)} chars")
                return text[:MAX_CHARS]
            try:
                buf.seek(0)
                import docx
                doc = docx.Document(buf)
                text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
                if text.strip():
                    logger.info(f"Parsed docx via python-docx: {len(text)} chars")
                    return text[:MAX_CHARS]
            except Exception as e:
                logger.warning(f"python-docx also failed: {e}")
            return None

        if "presentationml" in mime_type or mime_type == "application/vnd.ms-powerpoint":
            text = _parse_pptx(buf)
            if text:
                logger.info(f"Parsed pptx via stdlib: {len(text)} chars")
                return text[:MAX_CHARS]
            try:
                buf.seek(0)
                from pptx import Presentation
                prs = Presentation(buf)
                slides_text = []
                for i, slide in enumerate(prs.slides, 1):
                    parts = [f"--- Slide {i} ---"]
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                if para.text.strip():
                                    parts.append(para.text)
                    slides_text.append("\n".join(parts))
                text = "\n\n".join(slides_text)
                if text.strip():
                    logger.info(f"Parsed pptx via python-pptx: {len(text)} chars")
                    return text[:MAX_CHARS]
            except Exception as e:
                logger.warning(f"python-pptx also failed: {e}")
            return None

    except Exception as e:
        logger.error(f"Failed to parse office file ({mime_type}): {e}", exc_info=True)
        return None

    logger.warning(f"No parser matched for mime type: {mime_type}")
    return None


def download_file_content(service, file: dict) -> tuple[Optional[str], bool]:
    """Returns (content, is_truncated)"""
    mime_type = file["mimeType"]
    file_id = file["id"]

    try:
        text = None
        export_mime = GOOGLE_EXPORT_MAP.get(mime_type)
        if export_mime:
            res = service.files().export(fileId=file_id, mimeType=export_mime).execute()
            text = res.decode("utf-8") if isinstance(res, bytes) else str(res)
        
        elif mime_type == "application/pdf":
            request = service.files().get_media(fileId=file_id)
            buf = io.BytesIO()
            downloader = MediaIoBaseDownload(buf, request)
            done = False
            while not done:
                _, done = downloader.next_chunk()
            buf.seek(0)
            try:
                import PyPDF2
                reader = PyPDF2.PdfReader(buf)
                text = "\n".join(page.extract_text() or "" for page in reader.pages)
            except Exception:
                return None, False

        # Uploaded Office files — download binary and parse
        elif mime_type.startswith("application/vnd.openxmlformats-officedocument") or \
           mime_type in ("application/vnd.ms-excel", "application/msword", "application/vnd.ms-powerpoint"):
            request = service.files().get_media(fileId=file_id)
            buf = io.BytesIO()
            downloader = MediaIoBaseDownload(buf, request)
            done = False
            while not done:
                _, done = downloader.next_chunk()
            buf.seek(0)
            text = _parse_office_file(buf, mime_type)
        
        else:
            request = service.files().get_media(fileId=file_id)
            buf = io.BytesIO()
            downloader = MediaIoBaseDownload(buf, request)
            done = False
            while not done:
                _, done = downloader.next_chunk()
            buf.seek(0)
            text = buf.read().decode("utf-8", errors="replace")

        if text is not None:
            is_truncated = len(text) > MAX_CHARS
            return text[:MAX_CHARS], is_truncated
        
        return None, False

    except Exception:
        return None, False
